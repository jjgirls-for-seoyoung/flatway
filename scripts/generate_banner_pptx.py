import os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_banner_pptx():
    prs = Presentation()
    # 1:2 Scale for standard 600mm x 1800mm banner (PowerPoint slide max limit is 56 inches)
    # 300mm x 900mm maintains exact 1:3 proportion and scales losslessly to 600mm x 1800mm for print
    prs.slide_width = Mm(300)
    prs.slide_height = Mm(900)
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Scale factor from 600x1800 to 300x900
    S = 0.5

    def s_mm(val):
        return Mm(val * S)

    def s_pt(val):
        return Pt(val * S * 1.6)

    # MongoDB Design Tokens
    C_BRAND_GREEN = RGBColor(0, 237, 100)       # #00ED64
    C_BRAND_GREEN_DARK = RGBColor(0, 104, 74)   # #00684A
    C_BRAND_GREEN_SOFT = RGBColor(195, 240, 210)# #C3F0D2
    C_BRAND_TEAL_DEEP = RGBColor(0, 30, 43)     # #001E2B
    C_BRAND_TEAL = RGBColor(0, 61, 79)          # #003D4F
    C_ON_PRIMARY = RGBColor(0, 30, 43)          # #001E2B
    C_ACCENT_ORANGE = RGBColor(250, 110, 57)    # #FA6E39
    C_ACCENT_PURPLE = RGBColor(123, 63, 242)    # #7B3FF2
    C_ACCENT_BLUE = RGBColor(61, 79, 159)       # #3D4F9F
    C_CANVAS = RGBColor(255, 255, 255)
    C_SURFACE = RGBColor(249, 251, 250)         # #F9FBFA
    C_SURFACE_FEATURE = RGBColor(227, 252, 239) # #E3FCEF
    C_HAIRLINE = RGBColor(225, 229, 232)        # #E1E5E8
    C_HAIRLINE_STRONG = RGBColor(193, 204, 214) # #C1CCD6
    C_INK = RGBColor(0, 30, 43)                 # #001E2B
    C_STEEL = RGBColor(92, 108, 122)            # #5C6C7A
    C_ON_DARK = RGBColor(255, 255, 255)
    C_ON_DARK_MUTED = RGBColor(168, 179, 188)   # #A8B3BC

    # Background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, s_mm(0), s_mm(0), s_mm(600), s_mm(1800))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = C_SURFACE
    bg_shape.line.fill.background()

    # -------------------------------------------------------------
    # HEADER AREA (0 ~ 175mm) - MongoDB hero-band-dark
    # -------------------------------------------------------------
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, s_mm(0), s_mm(0), s_mm(600), s_mm(175))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = C_BRAND_TEAL_DEEP
    header_box.line.fill.background()

    # Signature MongoDB Green bottom accent line
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, s_mm(0), s_mm(170), s_mm(600), s_mm(5))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = C_BRAND_GREEN
    accent_line.line.fill.background()

    # Category Pill Badge (MongoDB badge-green)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(130), s_mm(16), s_mm(340), s_mm(24))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_BRAND_GREEN
    badge.line.fill.background()
    tf_b = badge.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "AI · 디지털 기반 문제해결 · 창작 프로젝트"
    p_b.font.size = s_pt(13.5)
    p_b.font.bold = True
    p_b.font.color.rgb = C_ON_PRIMARY
    p_b.font.name = "맑은 고딕"
    p_b.alignment = PP_ALIGN.CENTER

    # Main Title
    tx_title = slide.shapes.add_textbox(s_mm(20), s_mm(42), s_mm(560), s_mm(45))
    tf_t = tx_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "플랫웨이 (FlatWay)"
    p_t.font.size = s_pt(40)
    p_t.font.bold = True
    p_t.font.color.rgb = C_ON_DARK
    p_t.font.name = "맑은 고딕"
    p_t.alignment = PP_ALIGN.CENTER

    # Subtitle
    tx_sub = slide.shapes.add_textbox(s_mm(20), s_mm(92), s_mm(560), s_mm(26))
    tf_s = tx_sub.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = "휠체어와 이동 약자를 위한 실시간 단차 및 노면 파손 안내 지도"
    p_s.font.size = s_pt(17)
    p_s.font.bold = True
    p_s.font.color.rgb = C_BRAND_GREEN_SOFT
    p_s.font.name = "맑은 고딕"
    p_s.alignment = PP_ALIGN.CENTER

    # Slogan Banner (MongoDB dark card with green outline)
    slogan = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(50), s_mm(124), s_mm(500), s_mm(34))
    slogan.fill.solid()
    slogan.fill.fore_color.rgb = C_BRAND_TEAL
    slogan.line.color.rgb = C_BRAND_GREEN
    slogan.line.width = Pt(1.5)
    tf_sl = slogan.text_frame
    tf_sl.word_wrap = True
    p_sl = tf_sl.paragraphs[0]
    p_sl.text = "“턱 없는 세상, 안심하고 달리는 평평한 길을 연결합니다”"
    p_sl.font.size = s_pt(15)
    p_sl.font.bold = True
    p_sl.font.color.rgb = C_ON_DARK
    p_sl.font.name = "맑은 고딕"
    p_sl.alignment = PP_ALIGN.CENTER

    # Helper function for Section Headers (MongoDB style)
    def add_section_header(title_num, title_text, y_pos):
        h_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(y_pos), s_mm(550), s_mm(26))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = C_BRAND_TEAL_DEEP
        h_box.line.fill.background()
        
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(y_pos), s_mm(8), s_mm(26))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_BRAND_GREEN
        bar.line.fill.background()

        tf = h_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"  {title_num}. {title_text}"
        p.font.size = s_pt(15)
        p.font.bold = True
        p.font.color.rgb = C_ON_DARK
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------------------
    # 1. 기본 정보 (190 ~ 265mm) - MongoDB card-base
    # -------------------------------------------------------------
    add_section_header("1", "기본 정보", 190)
    
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(220), s_mm(550), s_mm(45))
    card1.fill.solid()
    card1.fill.fore_color.rgb = C_CANVAS
    card1.line.color.rgb = C_HAIRLINE
    card1.line.width = Pt(1)

    # 3 Pills
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(35), s_mm(226), s_mm(165), s_mm(33))
    b1.fill.solid()
    b1.fill.fore_color.rgb = C_SURFACE
    b1.line.color.rgb = C_HAIRLINE_STRONG
    tf1 = b1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "[학교명]"
    p1.font.size = s_pt(10.5)
    p1.font.bold = True
    p1.font.color.rgb = C_STEEL
    p1.font.name = "맑은 고딕"
    p1.alignment = PP_ALIGN.CENTER
    p1_2 = tf1.add_paragraph()
    p1_2.text = "작전여자고등학교"
    p1_2.font.size = s_pt(13)
    p1_2.font.bold = True
    p1_2.font.color.rgb = C_INK
    p1_2.font.name = "맑은 고딕"
    p1_2.alignment = PP_ALIGN.CENTER

    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(217), s_mm(226), s_mm(165), s_mm(33))
    b2.fill.solid()
    b2.fill.fore_color.rgb = C_SURFACE_FEATURE
    b2.line.color.rgb = C_BRAND_GREEN
    tf2 = b2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "[팀명]"
    p2.font.size = s_pt(10.5)
    p2.font.bold = True
    p2.font.color.rgb = C_BRAND_GREEN_DARK
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER
    p2_2 = tf2.add_paragraph()
    p2_2.text = "사대천왕"
    p2_2.font.size = s_pt(14)
    p2_2.font.bold = True
    p2_2.font.color.rgb = C_BRAND_GREEN_DARK
    p2_2.font.name = "맑은 고딕"
    p2_2.alignment = PP_ALIGN.CENTER

    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(400), s_mm(226), s_mm(165), s_mm(33))
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(250, 245, 255)
    b3.line.color.rgb = RGBColor(233, 213, 255)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "[프로젝트명]"
    p3.font.size = s_pt(10.5)
    p3.font.bold = True
    p3.font.color.rgb = C_ACCENT_PURPLE
    p3.font.name = "맑은 고딕"
    p3.alignment = PP_ALIGN.CENTER
    p3_2 = tf3.add_paragraph()
    p3_2.text = "플랫웨이 (FlatWay)"
    p3_2.font.size = s_pt(13)
    p3_2.font.bold = True
    p3_2.font.color.rgb = C_ACCENT_PURPLE
    p3_2.font.name = "맑은 고딕"
    p3_2.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # 2. 문제 정의 (275 ~ 445mm) - MongoDB card-base
    # -------------------------------------------------------------
    add_section_header("2", "문제 정의 (배경 및 필요성)", 275)

    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(305), s_mm(550), s_mm(140))
    card2.fill.solid()
    card2.fill.fore_color.rgb = C_CANVAS
    card2.line.color.rgb = C_HAIRLINE
    card2.line.width = Pt(1)

    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = s_mm(14)
    tf_c2.margin_right = s_mm(14)
    tf_c2.margin_top = s_mm(8)

    p = tf_c2.paragraphs[0]
    p.text = "● [발견한 문제]"
    p.font.size = s_pt(13)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_ORANGE
    p.font.name = "맑은 고딕"

    p = tf_c2.add_paragraph()
    p.text = "• 휠체어·유모차 이용자는 보행로의 높은 단차(턱)나 파손된 노면을 미리 알 수 없어, 이미 진입한 길에서 되돌아오거나 위험한 차도로 우회해야 하는 심각한 이동권 침해를 겪음."
    p.font.size = s_pt(11.5)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    p = tf_c2.add_paragraph()
    p.text = "● [해결 필요성 (기존 한계)]"
    p.font.size = s_pt(13)
    p.font.bold = True
    p.font.color.rgb = C_BRAND_GREEN_DARK
    p.font.name = "맑은 고딕"

    p = tf_c2.add_paragraph()
    p.text = "• 기존 상용 지도 앱은 차량 위주이거나 단순 최단 도보 경로만 제공하여 휠체어 통행에 필수적인 노면 파손·단차·경사도 등 세부 보행 장애 정보를 전혀 제공하지 못함."
    p.font.size = s_pt(11.5)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    # Persona Quote (MongoDB why-card)
    pq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(398), s_mm(524), Mm(22))
    pq.fill.solid()
    pq.fill.fore_color.rgb = C_SURFACE_FEATURE
    pq.line.color.rgb = C_BRAND_GREEN_SOFT
    tf_pq = pq.text_frame
    tf_pq.word_wrap = True
    tf_pq.margin_left = s_mm(8)
    tf_pq.margin_right = s_mm(8)
    tf_pq.margin_top = s_mm(3)
    p = tf_pq.paragraphs[0]
    p.text = "● 페르소나 (김민준, 24세 전동휠체어 대학생): “등하굣길에 턱에 막혀 차도로 돌아가야 했던 위험한 순간이 많았습니다. 미리 알고 피할 수 있다면 안심하고 다닐 수 있어요.”"
    p.font.size = s_pt(10.5)
    p.font.color.rgb = C_BRAND_GREEN_DARK
    p.font.name = "맑은 고딕"

    # -------------------------------------------------------------
    # 3. 해결 아이디어 (455 ~ 630mm) - MongoDB 3-tier
    # -------------------------------------------------------------
    add_section_header("3", "해결 아이디어 (핵심 솔루션)", 455)

    card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(485), s_mm(550), s_mm(140))
    card3.fill.solid()
    card3.fill.fore_color.rgb = C_CANVAS
    card3.line.color.rgb = C_HAIRLINE
    card3.line.width = Pt(1)

    sol_items = [
        ("① 맞춤 이동 수단별 안전 경로", "보행자 / 수동휠체어 / 전동휠체어 특성을 반영하여 '단차 최소 경로' 및 '파손 최소 경로'를 구분하여 최적의 안전 우회로 추천", C_SURFACE_FEATURE, C_BRAND_GREEN, C_BRAND_GREEN_DARK),
        ("② 실시간 참여형 위험 제보", "보행 중 발견한 턱이나 파손 도로를 사용자가 사진과 함께 원클릭으로 등록하고 지도에 실시간 공유", RGBColor(250, 245, 255), RGBColor(233, 213, 255), C_ACCENT_PURPLE),
        ("③ 건물 접근성 정보 연계", "도착지 건물의 주출입구 경사로 유무 및 엘리베이터 위치 정보를 상세히 제공하여 헛걸음 방지", RGBColor(239, 246, 255), RGBColor(191, 219, 254), C_ACCENT_BLUE)
    ]

    for idx, (s_title, s_desc, fill_c, line_c, text_c) in enumerate(sol_items):
        s_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(495 + idx * 42), s_mm(524), s_mm(38))
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = fill_c
        s_box.line.color.rgb = line_c
        s_box.line.width = Pt(1)
        tf_s = s_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = s_mm(8)
        tf_s.margin_right = s_mm(8)
        p = tf_s.paragraphs[0]
        p.text = s_title
        p.font.size = s_pt(12.5)
        p.font.bold = True
        p.font.color.rgb = text_c
        p.font.name = "맑은 고딕"
        p = tf_s.add_paragraph()
        p.text = s_desc
        p.font.size = s_pt(11)
        p.font.color.rgb = C_INK
        p.font.name = "맑은 고딕"

    # -------------------------------------------------------------
    # 4. AI·디지털 활용 과정 (635 ~ 805mm) - MongoDB SIGNATURE DARK TERMINAL CARD
    # -------------------------------------------------------------
    add_section_header("4", "AI · 디지털 도구 활용 과정", 635)

    # Dark Card (code-mockup-card)
    card4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(665), s_mm(550), s_mm(135))
    card4.fill.solid()
    card4.fill.fore_color.rgb = C_BRAND_TEAL_DEEP
    card4.line.color.rgb = C_BRAND_TEAL
    card4.line.width = Pt(1.5)

    tech_items = [
        ("스마트폰 센서 (IoT & Data)", "GPS 및 가속도 센서 연동 알고리즘", "휠체어 이동 중 발생하는 진동·충격 가속도 데이터를 수집하여 도로 노면의 파손 정도와 단차 위치를 자동 감지", C_BRAND_GREEN),
        ("지도 API 및 공간정보", "Kakao Maps / VWorld API 연동", "실제 도로 네트워크 및 보행로 데이터에 경사도와 노면 상태 레이어를 융합하여 휠체어 전용 가중치 기반 최적 경로 탐색 로직 설계", RGBColor(56, 189, 248)),
        ("UI/UX 디자인 도구", "Figma & Euclid 컴포넌트 설계", "이동 약자의 사용 편의성과 즉각적인 인지를 위해 '검색 → 수단 선택 → 맞춤 안내' 3단계 직관적 사용자 인터페이스 목업 설계", C_BRAND_GREEN_SOFT)
    ]

    for idx, (t_title, t_sub, t_desc, code_c) in enumerate(tech_items):
        t_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(675 + idx * 40), s_mm(524), s_mm(36))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = C_BRAND_TEAL
        t_box.line.color.rgb = RGBColor(28, 45, 56)
        t_box.line.width = Pt(1)
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = s_mm(8)
        tf_t.margin_right = s_mm(8)
        p = tf_t.paragraphs[0]
        p.text = f"▶ {t_title}   |   {t_sub}"
        p.font.size = s_pt(12)
        p.font.bold = True
        p.font.color.rgb = code_c
        p.font.name = "맑은 고딕"
        p = tf_t.add_paragraph()
        p.text = t_desc
        p.font.size = s_pt(10.5)
        p.font.color.rgb = C_ON_DARK
        p.font.name = "맑은 고딕"

    # -------------------------------------------------------------
    # 5. 테스트 및 개선 과정 (810 ~ 980mm) - MongoDB card-base
    # -------------------------------------------------------------
    add_section_header("5", "테스트 및 개선 과정", 810)

    card5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(840), s_mm(550), s_mm(135))
    card5.fill.solid()
    card5.fill.fore_color.rgb = C_CANVAS
    card5.line.color.rgb = C_HAIRLINE
    card5.line.width = Pt(1)

    # Sub box 1 (Mint Feature)
    test1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(850), s_mm(524), s_mm(56))
    test1.fill.solid()
    test1.fill.fore_color.rgb = C_SURFACE_FEATURE
    test1.line.color.rgb = C_BRAND_GREEN
    tf_t1 = test1.text_frame
    tf_t1.word_wrap = True
    tf_t1.margin_left = s_mm(8)
    tf_t1.margin_right = s_mm(8)
    p = tf_t1.paragraphs[0]
    p.text = "● [현장 실증 답사] 작전여고 ~ 작전역 구간 보행로 현장 조사 및 데이터 검증"
    p.font.size = s_pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_BRAND_GREEN_DARK
    p.font.name = "맑은 고딕"
    p = tf_t1.add_paragraph()
    p.text = "• 학교 주변 실제 통행로를 직접 답사하며 보도턱 높이(3~7cm) 및 노면 파손 지점을 실측 조사함.\n• 수집된 실측 데이터와 지도 API 경로를 비교 분석하여 휠체어 통행 불가 지점을 정확히 식별함."
    p.font.size = s_pt(10.8)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    # Sub box 2
    test2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(910), s_mm(524), s_mm(56))
    test2.fill.solid()
    test2.fill.fore_color.rgb = RGBColor(239, 246, 255)
    test2.line.color.rgb = RGBColor(191, 219, 254)
    tf_t2 = test2.text_frame
    tf_t2.word_wrap = True
    tf_t2.margin_left = s_mm(8)
    tf_t2.margin_right = s_mm(8)
    p = tf_t2.paragraphs[0]
    p.text = "● [개선 내용] 사용자 제보 신뢰도 검증 및 도로 상태 최신화 프로세스 설계"
    p.font.size = s_pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_BLUE
    p.font.name = "맑은 고딕"
    p = tf_t2.add_paragraph()
    p.text = "• 초기 아이디어의 허위/오래된 제보 문제를 해결하기 위해 '운영자 1차 필터링 + 다수 사용자 교차 확인' 시스템을 추가 설계하여 데이터의 신뢰성과 실시간성을 대폭 개선함."
    p.font.size = s_pt(10.8)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    # -------------------------------------------------------------
    # 6. 결과물 주요 화면 (985 ~ 1445mm)
    # -------------------------------------------------------------
    add_section_header("6", "결과물 주요 화면 및 핵심 기능", 985)

    card6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(1015), s_mm(550), s_mm(420))
    card6.fill.solid()
    card6.fill.fore_color.rgb = C_CANVAS
    card6.line.color.rgb = C_HAIRLINE
    card6.line.width = Pt(1)

    mockup_files = [
        ("assets/mockup_screen1.png", 38, "화면 ① [출발·도착지 & 단차 제보]", "• 현재 위치 기반 목적지 검색\n• 현장 사진 첨부 실시간 제보 UI"),
        ("assets/mockup_screen2.png", 218, "화면 ② [이동 수단 맞춤 선택]", "• 보행자 / 수동·전동 휠체어 선택\n• 수단별 예상시간 & 안전도 비교"),
        ("assets/mockup_screen3.png", 397, "화면 ③ [맞춤 안전경로 안내]", "• 파손 최소 안심길 & 단차 우회 안내\n• 경사로/엘리베이터 출입구 표시")
    ]

    for m_path, x_pos, m_title, m_desc in mockup_files:
        if os.path.exists(m_path):
            slide.shapes.add_picture(m_path, s_mm(x_pos), s_mm(1028), width=s_mm(165), height=s_mm(330))
        
        cap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(x_pos), s_mm(1365), s_mm(165), s_mm(60))
        cap_box.fill.solid()
        cap_box.fill.fore_color.rgb = C_SURFACE
        cap_box.line.color.rgb = C_HAIRLINE
        tf_cap = cap_box.text_frame
        tf_cap.word_wrap = True
        tf_cap.margin_left = s_mm(4)
        tf_cap.margin_right = s_mm(4)
        tf_cap.margin_top = s_mm(4)
        p = tf_cap.paragraphs[0]
        p.text = m_title
        p.font.size = s_pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_BRAND_TEAL_DEEP
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p = tf_cap.add_paragraph()
        p.text = m_desc
        p.font.size = s_pt(9.8)
        p.font.color.rgb = C_INK
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------------------
    # 7. 기대효과 (1450 ~ 1615mm)
    # -------------------------------------------------------------
    add_section_header("7", "기대효과 및 활용 가능성", 1450)

    card7 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(1480), s_mm(550), s_mm(130))
    card7.fill.solid()
    card7.fill.fore_color.rgb = C_CANVAS
    card7.line.color.rgb = C_HAIRLINE
    card7.line.width = Pt(1)

    imp1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(38), s_mm(1492), s_mm(255), s_mm(108))
    imp1.fill.solid()
    imp1.fill.fore_color.rgb = C_SURFACE_FEATURE
    imp1.line.color.rgb = C_BRAND_GREEN
    tf_i1 = imp1.text_frame
    tf_i1.word_wrap = True
    tf_i1.margin_left = s_mm(8)
    tf_i1.margin_right = s_mm(8)
    p = tf_i1.paragraphs[0]
    p.text = "● 이동 약자의 안전한 통행권 보장"
    p.font.size = s_pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_BRAND_GREEN_DARK
    p.font.name = "맑은 고딕"
    p = tf_i1.add_paragraph()
    p.text = "• 도로 턱과 파손 구역을 사전에 회피하여 통행 막힘 및 위험 차도 우회 사고를 원천 방지함.\n• 불필요한 우회 이동 시간을 단축하고 외출 시의 심리적 불안감을 크게 해소함."
    p.font.size = s_pt(10.5)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    imp2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(307), s_mm(1492), s_mm(255), s_mm(108))
    imp2.fill.solid()
    imp2.fill.fore_color.rgb = RGBColor(250, 245, 255)
    imp2.line.color.rgb = RGBColor(233, 213, 255)
    tf_i2 = imp2.text_frame
    tf_i2.word_wrap = True
    tf_i2.margin_left = s_mm(8)
    tf_i2.margin_right = s_mm(8)
    p = tf_i2.paragraphs[0]
    p.text = "● 배리어프리 도시 인프라 연계"
    p.font.size = s_pt(12.5)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_PURPLE
    p.font.name = "맑은 고딕"
    p = tf_i2.add_paragraph()
    p.text = "• 시민 참여형 노면 데이터를 축적하여 지자체 도로 보수 행정 및 우선 정비 구역 선정에 활용 가능.\n• 유모차, 고령자 등 모든 보행 약자로 수혜를 확장하여 포용적 도시 환경 조성."
    p.font.size = s_pt(10.5)
    p.font.color.rgb = C_INK
    p.font.name = "맑은 고딕"

    # -------------------------------------------------------------
    # 8. QR코드 삽입 영역 (1625 ~ 1775mm)
    # -------------------------------------------------------------
    card8 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(25), s_mm(1625), s_mm(550), s_mm(145))
    card8.fill.solid()
    card8.fill.fore_color.rgb = C_SURFACE
    card8.line.color.rgb = C_HAIRLINE_STRONG
    card8.line.width = Pt(1.5)

    qr_sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_mm(255), s_mm(1638), s_mm(90), s_mm(90))
    qr_sq.fill.solid()
    qr_sq.fill.fore_color.rgb = C_CANVAS
    qr_sq.line.color.rgb = C_BRAND_GREEN
    qr_sq.line.width = Pt(1.5)
    tf_qrs = qr_sq.text_frame
    tf_qrs.word_wrap = True
    p = tf_qrs.paragraphs[0]
    p.text = "QR CODE\n삽입 영역"
    p.font.size = s_pt(12)
    p.font.bold = True
    p.font.color.rgb = C_BRAND_GREEN_DARK
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    tx_qr_n = slide.shapes.add_textbox(s_mm(30), s_mm(1736), s_mm(540), s_mm(26))
    tf_qn = tx_qr_n.text_frame
    tf_qn.word_wrap = True
    p = tf_qn.paragraphs[0]
    p.text = "※ QR코드는 운영사에서 추후 일괄 적용 예정"
    p.font.size = s_pt(12)
    p.font.bold = True
    p.font.color.rgb = C_STEEL
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    # Save PPTX
    out_pptx = "flatway_xbanner_600x1800.pptx"
    prs.save(out_pptx)
    print(f"MongoDB PPTX saved successfully to {out_pptx}")

if __name__ == '__main__':
    create_banner_pptx()
